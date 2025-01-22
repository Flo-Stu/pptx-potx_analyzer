from flask import Flask, request, render_template, jsonify, send_file, Response, session
from pptx import Presentation
import os
from werkzeug.utils import secure_filename
import uuid
import json
from io import StringIO
from typing import Dict, Any
import time
import re
import html
from pptx.dml.chart import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import logging
from logging.handlers import RotatingFileHandler
from os import environ

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads/'
app.config['SECRET_KEY'] = environ.get('SECRET_KEY') or 'dev-key-please-change'
app.config['MAX_CONTENT_LENGTH'] = int(environ.get('MAX_CONTENT_LENGTH', 64 * 1024 * 1024))
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Implementierung eines Cleanup-Systems für temporäre Dateien
generated_files = {}  # Speichert Dateien und deren Erstellungszeit
TTL = 7 * 24 * 60 * 60  # Lebensdauer der Datei (7 Tage)

if not app.debug:
    if not os.path.exists('logs'):
        os.mkdir('logs')
    file_handler = RotatingFileHandler('logs/pptx_analyzer.log', maxBytes=10240, backupCount=10)
    file_handler.setFormatter(logging.Formatter(
        '%(asctime)s %(levelname)s: %(message)s [in %(pathname)s:%(lineno)d]'
    ))
    file_handler.setLevel(logging.INFO)
    app.logger.addHandler(file_handler)
    app.logger.setLevel(logging.INFO)
    app.logger.info('PPTX Analyzer startup')

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        if 'file' not in request.files:
            return 'Keine Datei ausgewählt'
        file = request.files['file']
        if file.filename == '':
            return 'Keine Datei ausgewählt'
        if file and (file.filename.endswith('.pptx') or file.filename.endswith('.potx')):
            # Sicherer Dateiname mit UUID
            filename = secure_filename(str(uuid.uuid4()) + '_' + file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)
            try:
                details = analyze_pptx(filepath)
                os.remove(filepath)  # Datei nach Analyse löschen
                return render_template('index.html', details=details)
            except Exception as e:
                os.remove(filepath)  # Datei auch bei Fehler löschen
                return jsonify({
                    'error': str(e),
                    'status': 'error',
                    'details': {
                        'type': type(e).__name__,
                        'message': str(e)
                    }
                }), 500
    return render_template('index.html', details=None)

def analyze_pptx(filepath):
    prs = Presentation(filepath)
    details = []
    for i, layout in enumerate(prs.slide_master.slide_layouts):
        layout_info = {
            'layout_index': i,
            'layout_name': layout.name,
            'placeholders': []
        }
        for placeholder in layout.placeholders:
            layout_info['placeholders'].append({
                'idx': placeholder.placeholder_format.idx,
                'type': placeholder.placeholder_format.type
            })
        details.append(layout_info)
    return details

@app.route('/generate_script', methods=['POST'])
def generate_script():
    master_url = request.form.get('master_url')
    layout_data = json.loads(request.form.get('layout_data'))
    
    # Generiere den Python-Code
    generated_code = generate_pptx_code(master_url, layout_data)
    session['last_generated_code'] = generated_code
    
    return render_template('index.html', details=layout_data, generated_code=generated_code)

@app.route('/download_script')
def download_script():
    generated_code = session.get('last_generated_code', '')
    if not generated_code:
        return "Kein Code verfügbar", 404
    
    return Response(
        generated_code,
        mimetype='text/plain',
        headers={'Content-Disposition': 'attachment;filename=create_presentation.py'}
    )

def generate_pptx_code(master_url, layout_data):
    code = f"""from pptx import Presentation
import requests
from io import BytesIO
from typing import Dict, Any

# Definiere die verfügbaren Layouts aus der Analyse
AVAILABLE_LAYOUTS = {json.dumps(layout_data, indent=4)}

def create_presentation(data: Dict[str, Any]) -> None:
    # Lade Master-Template von URL
    response = requests.get("{master_url}")
    prs = Presentation(BytesIO(response.content))
    
    def create_slide(layout_index: int, content_dict: Dict[int, str]) -> None:
        \"\"\"
        Erstellt eine neue Folie mit dem spezifizierten Layout und Inhalt.
        
        Args:
            layout_index: Index des zu verwendenden Layouts (0-X, basierend auf der Analyse)
            content_dict: Dictionary mit den verfügbaren Platzhalter-IDs aus der Analyse
        \"\"\"
        if layout_index >= len(prs.slide_layouts):
            raise ValueError(f"Layout-Index {{layout_index}} nicht verfügbar")
            
        layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(layout)
        
        # Fülle die Platzhalter
        for shape in slide.placeholders:
            idx = shape.placeholder_format.idx
            if idx in content_dict:
                shape.text = content_dict[idx]

    # Erstelle die Folien basierend auf den Daten
    for slide_data in data["slides"]:
        create_slide(
            layout_index=slide_data["layout_index"],
            content_dict=slide_data["content"]
        )
    
    # Speichere die Präsentation
    output_path = data.get("output_path", "neue_praesentation.pptx")
    prs.save(output_path)
"""
    return code

def escape_text(text: str) -> str:
    """Entfernt Markdown und HTML-Artefakte."""
    if not text:
        return ""
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    text = html.unescape(text)
    return text

@app.errorhandler(413)
def too_large(e):
    return "Datei ist zu groß. Maximum ist 64MB", 413

@app.errorhandler(500)
def server_error(e):
    return "Interner Serverfehler. Bitte später erneut versuchen.", 500

@app.after_request
def add_security_headers(response):
    response.headers['X-Content-Type-Options'] = 'nosniff'
    response.headers['X-Frame-Options'] = 'SAMEORIGIN'
    response.headers['X-XSS-Protection'] = '1; mode=block'
    return response

if __name__ == '__main__':
    app.run(debug=False)